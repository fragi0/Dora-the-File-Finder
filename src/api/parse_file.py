from docx import Document
from PIL import Image
import pytesseract
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook
from pypdf import PdfReader
from fastapi import UploadFile
import zipfile
import xml.etree.ElementTree as ET


async def parse_file(file_like: UploadFile) -> str:
    if not file_like.filename or "." not in file_like.filename:
        raise IOError("Couldn't infer filetype, file with no extension.")

    data: str = ""
    ft = file_like.filename.rsplit(".", 1)[-1].lower()

    try:
        # ---------------- PDF ----------------
        if ft == "pdf":
            reader = PdfReader(file_like.file)
            for pag in range(len(reader.pages)):
                page = reader.pages[pag]
                text = page.extract_text()
                if text:
                    data += text + "\n"
            file_like.file.seek(0)

        # ---------------- DOCX ----------------
        elif ft == "docx":
            document = Document(file_like.file)
            paragraphs = [p.text for p in document.paragraphs
                          if p.text.strip()]
            paragraphs.extend(
                "\t".join(cell.text for cell in row.cells)
                for table in document.tables
                for row in table.rows)
            data = "\n".join(paragraphs)
            file_like.file.seek(0)

        # ---------------- IMAGES ----------------
        elif ft in ["png", "jpg", "jpeg"]:
            image = Image.open(file_like.file)
            max_size = (1024, 1024)
            image.thumbnail(max_size, Image.Resampling.LANCZOS)
            data = pytesseract.image_to_string(image).strip()
            file_like.file.seek(0)

        # ---------------- PPTX ----------------
        elif ft == "pptx":
            data = extract_text_from_pptx(file_like.file)
            file_like.file.seek(0)

        # ---------------- EXCEL ----------------
        elif ft in ["xls", "xlsx"]:
            wb = load_workbook(file_like.file, read_only=True)
            dfs = [
                pd.DataFrame(rows[1:], columns=rows[0])
                for sheetname in wb.sheetnames
                for rows in [list(wb[sheetname].values)]
                if rows
            ]
            if dfs:
                data = pd.concat(dfs).to_csv(index=False)
            else:
                data = ""
            file_like.file.seek(0)

        # ------- OPENDOCUMENT (ODT, ODS, ODP) -------
        elif ft in ["odt", "ods", "odp"]:
            data = extract_odf_text(file_like.file)
            file_like.file.seek(0)

        else:
            if await is_text_file(file_like):
                content = await file_like.read()
                data = content.decode('utf-8', errors='ignore')
                await file_like.seek(0)
            else:
                print(f"Error: {file_like.filename} binary files"
                      "are not supported.")

    except Exception as e:
        print(f"There was an error when processing the file: {e}")
        data = ""

    return data.strip()


def extract_text_from_pptx(file_like) -> str:
    presentation = Presentation(file_like)
    extracted_text: str = ""

    for slide_number, slide in enumerate(presentation.slides):
        extracted_text += f"\nSlide {slide_number + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                extracted_text += shape.text.strip() + "\n"

    return extracted_text.strip()


async def is_text_file(file: UploadFile) -> bool:
    try:
        chunk = await file.read(1024)
        await file.seek(0)
        if not chunk:
            return True
        if b'\x00' in chunk:
            return False
        chunk.decode('utf-8')
        return True
    except UnicodeDecodeError:
        return False


def extract_odf_text(file_like) -> str:
    text_parts = []
    try:
        with zipfile.ZipFile(file_like) as z:
            content_xml = z.read("content.xml")
            tree = ET.fromstring(content_xml)
            for node in tree.iter():
                if node.text and node.text.strip():
                    text_parts.append(node.text.strip())
        return "\n".join(text_parts)
    except Exception as e:
        print(f"OpenDocument file processing error: {e}")
        return ""
